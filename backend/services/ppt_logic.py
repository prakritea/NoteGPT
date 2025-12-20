import os
import json
import re
from typing import List, Dict

from pptx import Presentation
from backend.utils.gemini_helper import generate_gemini_response_async
from backend.models.schemas import SlideContent

# Base directory: <project root>/backend
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
# Absolute templates dir: <project root>/backend/templates
TEMPLATES_DIR = os.path.join(BASE_DIR, "templates")

# Friendly IDs → preferred filenames (these are just hints; resolver will search if needed)
TEMPLATE_MAP = {
    "template1": "Template 1.pptx",
    "template2": "Template 2.pptx",
    "template3": "Template 3.pptx",
    # add more mappings as needed
}


def list_available_templates() -> List[Dict[str, str]]:
  """
  Return a list of available PPT templates in backend/templates.

  Each item has:
  - id: filename (e.g. 'Template 1.pptx')
  - label: nice name without extension (e.g. 'Template 1')
  """
  if not os.path.isdir(TEMPLATES_DIR):
      return []

  templates: List[Dict[str, str]] = []
  for fname in os.listdir(TEMPLATES_DIR):
      if not fname.lower().endswith(".pptx"):
          continue
      label = os.path.splitext(fname)[0]
      templates.append({"id": fname, "label": label})

  # sort alphabetically by label for stable UI
  templates.sort(key=lambda t: t["label"].lower())
  return templates


def resolve_template_filename(template_id: str) -> str:
    """
    Try to find a matching filename under TEMPLATES_DIR for the given template_id.
    Accepts friendly ids like 'template1', 'template1.pptx', 'Template 1', etc.

    1) Try TEMPLATE_MAP
    2) Try direct filename / with .pptx
    3) Case-insensitive match in templates folder
    4) Substring match ignoring spaces/underscore/dash

    Raises FileNotFoundError if no match is found.
    """
    # 1) Try TEMPLATE_MAP
    mapped = TEMPLATE_MAP.get(template_id)
    if mapped:
        candidate_path = os.path.join(TEMPLATES_DIR, mapped)
        if os.path.exists(candidate_path):
            return mapped

    # 2) Normalize incoming name
    name = os.path.basename(template_id)
    if not name.lower().endswith(".pptx"):
        name_variants = [name, f"{name}.pptx", name.replace(" ", "_") + ".pptx"]
    else:
        name_variants = [name]

    if not os.path.isdir(TEMPLATES_DIR):
        raise FileNotFoundError(f"Templates directory not found at {TEMPLATES_DIR}")

    entries = os.listdir(TEMPLATES_DIR)

    # 3) Case-insensitive exact match
    for variant in name_variants:
        for e in entries:
            if e.lower() == variant.lower():
                return e

    # 4) Substring / fuzzy match (ignore spaces, underscores, dashes)
    tid_norm = template_id.lower().replace(" ", "").replace("_", "").replace("-", "")
    for e in entries:
        e_norm = e.lower().replace(" ", "").replace("_", "").replace("-", "")
        if tid_norm and tid_norm in e_norm:
            return e

    raise FileNotFoundError(
        f"No template file matching '{template_id}' found in {TEMPLATES_DIR}. "
        f"Available files: {entries}"
    )


def _extract_json_array(text: str) -> str:
    """
    Extract the first top-level JSON array from the model output.
    Handles cases where the model wraps JSON in markdown/code fences or extra text.
    """
    cleaned = re.sub(r"```(?:json)?", "", text, flags=re.IGNORECASE).strip()
    start = cleaned.find("[")
    end = cleaned.rfind("]")
    if start == -1 or end == -1 or end <= start:
        return cleaned
    return cleaned[start : end + 1]


async def generate_slide_content(topic: str, num_slides: int) -> List[SlideContent]:
    """
    Use Gemini to generate slide outlines. Returns a list of SlideContent objects.
    Each SlideContent has a title and a list of bullet strings.
    """
    prompt = f"""
You are an assistant that creates PowerPoint slide outlines.

Topic: {topic}
Generate exactly {num_slides} slides.

For EACH slide, provide:
- "title": a short, clear slide title
- "bullets": a list of 3-5 concise bullet points as strings

Return the result as a JSON array of objects, like this:
[
  {{
    "title": "Slide 1 title",
    "bullets": ["point 1", "point 2", "point 3"]
  }},
  {{
    "title": "Slide 2 title",
    "bullets": ["point 1", "point 2", "point 3"]
  }}
]

IMPORTANT:
- Return ONLY the JSON array.
- Do NOT include any explanations, markdown, or text before/after the JSON.
"""

    raw = await generate_gemini_response_async(prompt)
    if not raw:
        return []

    json_str = _extract_json_array(raw)

    try:
        slides_json = json.loads(json_str)
    except json.JSONDecodeError:
        # Very lenient fallback
        slides_json = []
        for line in raw.splitlines():
            line = line.strip()
            if not line or not line.startswith("{"):
                continue
            try:
                obj = json.loads(line.rstrip(","))
                slides_json.append(obj)
            except Exception:
                continue

    results: List[SlideContent] = []

    if isinstance(slides_json, list):
        for s in slides_json:
            if not isinstance(s, dict):
                continue
            title = str(s.get("title", "")).strip()
            bullets = s.get("bullets", [])
            if not title or not isinstance(bullets, list):
                continue

            cleaned_bullets = [
                str(b).strip()
                for b in bullets
                if isinstance(b, (str, int, float)) and str(b).strip()
            ]
            if cleaned_bullets:
                results.append(SlideContent(title=title, bullets=cleaned_bullets))

    return results


def _remove_slide(prs: Presentation, index: int) -> None:
    """
    Remove a slide at the given index from a Presentation.
    python-pptx has no official API for this, so we use the underlying XML list.
    """
    slide_id_list = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slide_ids = list(slide_id_list)
    slide_id_list.remove(slide_ids[index])


def build_ppt_from_content(
    slides: List[SlideContent], template_id: str, output_path: str
) -> str:
    """
    Build a PPTX by OVERWRITING the content of existing template slides.

    - Slide 0 (cover): only overwrite the title (no bullets).
    - Other slides: overwrite title + put bullets into a chosen text frame:
        1) Prefer BODY/SUBTITLE placeholders (if any).
        2) Prefer shapes on the right half of the slide (common for modern templates).
        3) Otherwise fall back to the largest text frame (by width * height).
    """
    template_filename = resolve_template_filename(template_id)
    template_path = os.path.join(TEMPLATES_DIR, template_filename)

    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template {template_filename} not found at {template_path}")

    prs = Presentation(template_path)

    try:
        default_layout = prs.slide_layouts[1]
    except Exception:
        default_layout = prs.slide_layouts[0]

    slide_width = prs.slide_width

    for idx, slide_content in enumerate(slides):
        if idx < len(prs.slides):
            slide = prs.slides[idx]
        else:
            slide = prs.slides.add_slide(default_layout)

        # ----- TITLE -----
        title_shape = None
        for shp in slide.shapes:
            if getattr(shp, "is_placeholder", False):
                phf = getattr(shp, "placeholder_format", None)
                if phf and phf.type in (1, 3):  # TITLE or CENTER_TITLE
                    title_shape = shp
                    break
        if title_shape is None and slide.shapes.title:
            title_shape = slide.shapes.title

        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            title_shape.text = slide_content.title

        # Cover slide: skip bullets
        if idx == 0:
            continue

        # ----- BODY -----
        candidates = []
        for shp in slide.shapes:
            if shp is title_shape:
                continue
            if not getattr(shp, "has_text_frame", False):
                continue
            candidates.append(shp)

        if not candidates:
            continue

        # (1) Prefer BODY / SUBTITLE placeholders
        placeholder_candidates = []
        for shp in candidates:
            if getattr(shp, "is_placeholder", False):
                phf = getattr(shp, "placeholder_format", None)
                if phf and phf.type in (2, 4):  # BODY, SUBTITLE
                    placeholder_candidates.append(shp)

        if placeholder_candidates:
            candidates = placeholder_candidates

        # (2) Prefer shapes on the right half of the slide
        right_side_candidates = []
        for shp in candidates:
            center_x = getattr(shp, "left", 0) + getattr(shp, "width", 0) // 2
            if center_x > slide_width // 2:
                right_side_candidates.append(shp)

        if right_side_candidates:
            candidates = right_side_candidates

        # (3) Choose the largest remaining candidate by area
        body_shape = None
        max_area = 0
        for shp in candidates:
            width = getattr(shp, "width", 0)
            height = getattr(shp, "height", 0)
            area = int(width) * int(height)
            if area > max_area:
                max_area = area
                body_shape = shp

        if body_shape is not None and getattr(body_shape, "has_text_frame", False):
            tf = body_shape.text_frame
            tf.clear()
            for b in slide_content.bullets:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0

    # Remove any extra template slides beyond what we used
    while len(prs.slides) > len(slides):
        _remove_slide(prs, len(prs.slides) - 1)

    prs.save(output_path)
    return output_path
